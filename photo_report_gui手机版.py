"""
华为门店家装验收自动化工具 - GUI 版本 v2.8
功能：
  - 单张模式：逐张审阅照片并生成单页PPT
  - 批量模式：一次填写多条记录，一键生成多页PPT
"""

import os
import sys
import json
import tkinter as tk
from tkinter import ttk, messagebox
from datetime import datetime
from PIL import Image, ImageTk, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy
import hashlib
import platform

# ====== 资源路径处理（支持打包后的EXE） ======
def get_resource_path(relative_path):
    """获取资源文件的绝对路径，支持PyInstaller打包"""
    try:
        base_path = sys._MEIPASS
    except AttributeError:
        base_path = os.path.dirname(os.path.abspath(__file__))
    return os.path.join(base_path, relative_path)

# ====== 配置区域 ======
BASE_DIR = os.path.dirname(os.path.abspath(sys.argv[0])) if getattr(sys, 'frozen', False) else os.path.dirname(os.path.abspath(__file__))

# 资源文件路径（支持打包后从 _MEIPASS 读取）
def _res(p):
    """获取资源路径，支持PyInstaller打包"""
    return get_resource_path(p)

WORK_DIR = BASE_DIR
PHOTO_FOLDER = os.path.join(BASE_DIR, "验收照片")
OUTPUT_FOLDER = os.path.join(BASE_DIR, "验收报告")
PLAN_IMAGES_FOLDER = os.path.join(BASE_DIR, "平面图")
TEMPLATE_PPTX = _res("单页模版.pptx")
DECOR_IMG = _res("装饰底图.jpg")  # 母版装饰图(MAKE it POSSIBLE+红色弧形)，嵌入每页保证复制不丢
SESSION_FILE = os.path.join(OUTPUT_FOLDER, "session_data.json")
SETTINGS_FILE = _res("settings.json")
ISSUES_DB_FILE = _res("issues_db.json")
REMEDIES_DB_FILE = _res("remedies_db.json")

LICENSE_FILE = _res("license.dat")
SECRET_KEY = "Huawei@2026!YongGe"

# ====== 跨平台字体 ======
def _get_cn_font():
    """返回当前平台可用的中文字体名"""
    _sys = platform.system()
    if _sys == 'Darwin':
        return 'PingFang SC'   # macOS 内置中文字体
    elif _sys == 'Windows':
        return '微软雅黑'
    else:
        return 'WenQuanYi Micro Hei'  # Linux 常见中文字体

CN_FONT = _get_cn_font()

def _open_file(path):
    """跨平台打开文件/文件夹"""
    try:
        _sys = platform.system()
        if _sys == 'Windows':
            os.startfile(path)
        elif _sys == 'Darwin':
            import subprocess
            subprocess.call(['open', path])
        else:
            import subprocess
            subprocess.call(['xdg-open', path])
    except Exception:
        pass

def get_machine_code():
    """获取机器唯一码 - 跨平台实现（Windows/Mac/Linux）"""
    import uuid
    parts = []

    # 方法1：Windows注册表 MachineGuid（仅Windows）
    if platform.system() == 'Windows':
        try:
            import winreg
            key = winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\Microsoft\Cryptography")
            machine_guid, _ = winreg.QueryValueEx(key, "MachineGuid")
            winreg.CloseKey(key)
            parts.append(machine_guid)
        except:
            pass

    # 方法2：Mac IOPlatformUUID（仅macOS）
    if platform.system() == 'Darwin':
        try:
            import subprocess
            result = subprocess.check_output(
                ['ioreg', '-rd1', '-c', 'IOPlatformExpertDevice'],
                stderr=subprocess.DEVNULL
            ).decode('utf-8', errors='ignore')
            for line in result.splitlines():
                if 'IOPlatformUUID' in line:
                    parts.append(line.split('"')[-2])
                    break
        except:
            pass

    # 方法3：计算机名（通用备用）
    try:
        parts.append(platform.node())
    except:
        pass

    # 方法4：网卡MAC地址（通用兜底）
    try:
        mac = ':'.join(f'{(uuid.getnode() >> i) & 0xff:02x}' for i in range(0, 48, 8)[::-1])
        parts.append(mac)
    except:
        pass

    raw = "|".join(filter(None, parts))
    return hashlib.sha256(raw.encode()).hexdigest()[:16]

def generate_license_key(machine_code):
    """生成授权密钥"""
    raw = machine_code + SECRET_KEY
    return hashlib.sha256(raw.encode()).hexdigest()[:24]

def check_license():
    """检查授权"""
    machine_code = get_machine_code()
    
    if not os.path.exists(LICENSE_FILE):
        return False, machine_code
    try:
        with open(LICENSE_FILE, 'r', encoding='utf-8') as f:
            saved_key = f.read().strip()
        valid_key = generate_license_key(machine_code)
        if saved_key == valid_key:
            return True, ''
    except Exception:
        pass
    return False, machine_code

def save_license(machine_code):
    """保存授权文件"""
    key = generate_license_key(machine_code)
    with open(LICENSE_FILE, 'w', encoding='utf-8') as f:
        f.write(key)
    return True

# 预设选项
CATEGORIES = ["店铺问题—基装", "店铺问题—家装"]

def _build_titles_from_db():
    """从 issues_db.json 动态读取所有标题，保证与数据库完全同步"""
    _default = {
        "店铺问题—基装": [
            "玻璃门/玻璃幕墙", "石材幕墙", "铝板幕墙", "石膏板吊顶", "铝板吊顶",
            "天花灯具/灯带", "软膜天花/其他平板灯具", "地砖/地面石材/地面岩板",
            "地毯类/木地板类", "石膏板墙面", "快装板/石材/吸音板/其他饰面材料",
            "木饰面/雕刻铝板", "烤漆门/暗门", "自动玻璃门/室内玻璃隔断", "楼梯",
            "墙地插座、开关", "配电箱/强电线路", "弱电设备", "空调设备", "消防设备"
        ],
        "店铺问题—家装": [
            "体验桌", "灯箱", "标准配件柜", "智慧屏柜", "门头", "LOGO", "导视标识"
        ]
    }
    if not os.path.exists(ISSUES_DB_FILE):
        return _default
    try:
        with open(ISSUES_DB_FILE, 'r', encoding='utf-8-sig') as _f:
            _db = json.load(_f)
        result = {}
        for _cat_key, _ui_key in [("基装", "店铺问题—基装"), ("家装", "店铺问题—家装")]:
            if _cat_key in _db:
                result[_ui_key] = list(_db[_cat_key].keys())
            else:
                result[_ui_key] = _default.get(_ui_key, [])
        return result
    except Exception:
        return _default

TITLES_BY_CATEGORY = _build_titles_from_db()
SEVERITIES = ["一般问题", "严重问题", "关键问题", "普通问题", "记录问题"]

# ====== 辅助函数 ======

def load_json(path, default=None):
    if os.path.exists(path):
        with open(path, 'r', encoding='utf-8-sig') as f:
            return json.load(f)
    return default if default is not None else []

def save_json(path, data):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

def load_db(path):
    data = load_json(path, {})
    return data if data else {}

def load_settings():
    if os.path.exists(SETTINGS_FILE):
        try:
            with open(SETTINGS_FILE, 'r', encoding='utf-8-sig') as f:
                return json.load(f)
        except:
            pass
    return {}

def save_settings(settings):
    with open(SETTINGS_FILE, 'w', encoding='utf-8') as f:
        json.dump(settings, f, ensure_ascii=False, indent=2)

# ====== 批量录入槽位控件 ======


class BatchSlot(ttk.Frame):
    """批量录入中的单个问题槽位 - 下方展开式照片选择面板"""

    THUMB_W = 180   # 增加缩略图宽度
    THUMB_H = 140   # 增加缩略图高度
    COLS    = 4

    def __init__(self, parent, index, app, **kwargs):
        super().__init__(parent, relief='groove', borderwidth=1, **kwargs)
        self.index = index
        self.app = app
        self.photo_path = None
        self.photo_name = None
        self._thumb = None
        self._preview_original_path = None
        self._selected_fname_var = tk.StringVar()   # 当前选中的文件名
        self._build()

    def _build(self):
        # 上部：基本信息区域（横向排列）
        top = ttk.Frame(self)
        top.pack(fill=tk.X, padx=6, pady=4)

        # 左侧：照片缩略图
        left = ttk.Frame(top)
        left.pack(side=tk.LEFT, padx=6, pady=6)

        ttk.Label(left, text=f"#{self.index + 1}", font=(CN_FONT, 9, 'bold')).pack()

        self.thumb_label = ttk.Label(left, text="未选照片", width=16,
                                     relief='sunken', anchor='center',
                                     background='#e0e0e0')
        self.thumb_label.pack(pady=3)
        self.thumb_label.config(cursor='hand2')
        self.thumb_label.bind('<Button-1>', lambda e: self._open_selected_photo())

        self.pick_btn = ttk.Button(left, text="选择照片", width=10,
                                   command=self._toggle_thumb_panel)
        self.pick_btn.pack(pady=2)

        # 右侧：表单（变窄以腾出空间）
        right = ttk.Frame(top)
        right.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=4, pady=6)

        r = 0
        ttk.Label(right, text="类别：", font=(CN_FONT, 9)).grid(row=r, column=0, sticky=tk.W, pady=2)
        self.cat_var = tk.StringVar(value=CATEGORIES[0])
        cat_cb = ttk.Combobox(right, textvariable=self.cat_var,
                               values=CATEGORIES, state='readonly', width=12, font=(CN_FONT, 9))
        cat_cb.grid(row=r, column=1, sticky=tk.W, pady=2)
        cat_cb.bind('<<ComboboxSelected>>', self._on_cat_changed)
        r += 1

        ttk.Label(right, text="标题：", font=(CN_FONT, 9)).grid(row=r, column=0, sticky=tk.W, pady=2)
        titles = TITLES_BY_CATEGORY.get(CATEGORIES[0], [])
        self.title_var = tk.StringVar(value=titles[0] if titles else '')
        self.title_cb = ttk.Combobox(right, textvariable=self.title_var,
                                     values=titles, width=12, font=(CN_FONT, 9))
        self.title_cb.grid(row=r, column=1, sticky=tk.W, pady=2)
        self.title_cb.bind('<Button-3>', self._title_right_click)
        r += 1

        ttk.Label(right, text="严重程度：", font=(CN_FONT, 9)).grid(row=r, column=0, sticky=tk.W, pady=2)
        self.sev_var = tk.StringVar(value=SEVERITIES[0])
        ttk.Combobox(right, textvariable=self.sev_var,
                     values=SEVERITIES, state='readonly', width=12, font=(CN_FONT, 9)).grid(
            row=r, column=1, sticky=tk.W, pady=2)
        r += 1

        ttk.Label(right, text="问题描述：", font=(CN_FONT, 9)).grid(row=r, column=0, sticky=tk.NW, pady=2)
        desc_frame = ttk.Frame(right)
        desc_frame.grid(row=r, column=1, sticky=tk.W, pady=2)
        self.desc_text = tk.Text(desc_frame, width=20, height=3, font=(CN_FONT, 9))
        self.desc_text.pack(side=tk.LEFT)
        ttk.Button(desc_frame, text="从库选", width=6,
                   command=self._pick_desc).pack(side=tk.LEFT, padx=3)
        r += 1

        ttk.Label(right, text="整改意见：", font=(CN_FONT, 9)).grid(row=r, column=0, sticky=tk.NW, pady=2)
        remedy_frame = ttk.Frame(right)
        remedy_frame.grid(row=r, column=1, sticky=tk.W, pady=2)
        self.remedy_text = tk.Text(remedy_frame, width=20, height=3, font=(CN_FONT, 9))
        self.remedy_text.pack(side=tk.LEFT)
        ttk.Button(remedy_frame, text="从库选", width=6,
                   command=self._pick_remedy).pack(side=tk.LEFT, padx=3)
        r += 1

        ttk.Label(right, text="平面图：", font=(CN_FONT, 9)).grid(row=r, column=0, sticky=tk.W, pady=2)
        self.plan_var = tk.StringVar(value="无")
        ttk.Combobox(right, textvariable=self.plan_var,
                     values=self.app.plan_images, state='readonly',
                     width=12, font=(CN_FONT, 9)).grid(row=r, column=1, sticky=tk.W, pady=2)

        # 可展开的缩略图面板（默认隐藏，在下方全宽展开）
        self._thumb_panel = ttk.Frame(self, relief='sunken', borderwidth=1)
        self._build_thumb_panel()

    def _toggle_thumb_panel(self):
        """展开或收起照片选择面板"""
        if self._thumb_panel.winfo_manager():
            self._thumb_panel.pack_forget()
            self.pick_btn.config(text="选择照片")
        else:
            self._thumb_panel.pack(fill=tk.BOTH, expand=True, padx=6, pady=(0, 6))
            self.pick_btn.config(text="收起照片")
            self._load_thumb_grid()
            if self.photo_name:
                self.after(150, self._auto_select)

    def _build_thumb_panel(self):
        """构建照片选择面板（在批次下方展开）"""
        panel = self._thumb_panel

        toolbar = ttk.Frame(panel, padding=(8, 6))
        toolbar.pack(fill=tk.X)

        ttk.Button(toolbar, text="浏览...",
                    command=self._browse_photo).pack(side=tk.LEFT, padx=(0, 6))

        self._status_var = tk.StringVar(value="加载中...")
        ttk.Label(toolbar, textvariable=self._status_var,
                   font=(CN_FONT, 8), foreground='#666').pack(side=tk.LEFT, padx=12)

        # 关闭按钮
        ttk.Button(toolbar, text="✕ 关闭", width=10,
                    command=self._toggle_thumb_panel).pack(side=tk.RIGHT)

        main = ttk.Frame(panel, padding=(8, 4))
        main.pack(fill=tk.BOTH, expand=True)

        # 左侧：缩略图网格（全宽）
        left = ttk.Frame(main)
        left.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=(0, 10))

        ttk.Label(left, text="照片缩略图（单击预览，双击直接选择）",
                   font=(CN_FONT, 9, 'bold')).pack(anchor='w')

        grid_outer = ttk.Frame(left)
        grid_outer.pack(fill=tk.BOTH, expand=True)

        vsb = ttk.Scrollbar(grid_outer, orient=tk.VERTICAL)
        vsb.pack(side=tk.RIGHT, fill=tk.Y)

        canvas = tk.Canvas(grid_outer, yscrollcommand=vsb.set, highlightthickness=0)
        canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        vsb.config(command=canvas.yview)

        inner = ttk.Frame(canvas)
        canvas_window = canvas.create_window((0, 0), window=inner, anchor='nw')
        self._pick_canvas = canvas
        self._pick_inner = inner
        self._pick_canvas_window = canvas_window

        def _on_canvas_resize(e):
            canvas.itemconfig(canvas_window, width=e.width)
        canvas.bind('<Configure>', _on_canvas_resize)

        def _on_inner_resize(e):
            canvas.configure(scrollregion=canvas.bbox("all"))
        inner.bind('<Configure>', _on_inner_resize)

        # 鼠标滚轮滚动
        def _on_mousewheel(e):
            canvas.yview_scroll(-1 * (e.delta // 120), "units")
        canvas.bind('<Enter>', lambda e: canvas.bind_all('<MouseWheel>', _on_mousewheel))
        canvas.bind('<Leave>', lambda e: canvas.unbind_all('<MouseWheel>'))

        # 键盘上下键滚动
        def _on_key_up(e):
            canvas.yview_scroll(-3, "units")
        def _on_key_down(e):
            canvas.yview_scroll(3, "units")
        canvas.bind('<Up>', _on_key_up)
        canvas.bind('<Down>', _on_key_down)
        canvas.config(takefocus=True)  # 允许获取焦点以接收键盘事件

        # 右侧：预览区（减小面积）
        right = ttk.Frame(main)
        right.pack(side=tk.LEFT, fill=tk.Y)

        ttk.Label(right, text="预览（双击查看原图）", font=(CN_FONT, 9, 'bold')).pack(anchor='w')

        preview_outer = ttk.Frame(right, relief='sunken', borderwidth=1,
                                  width=180, height=140)  # 减小预览区
        preview_outer.pack(pady=(4, 0))
        preview_outer.pack_propagate(False)

        self._preview_label = ttk.Label(preview_outer, text="请点击左侧缩略图",
                                          anchor='center', justify='center',
                                          cursor='hand2')
        self._preview_label.pack(fill=tk.BOTH, expand=True)
        self._preview_label.bind('<Double-Button-1>',
                                 lambda e: self._open_original_photo())

        self._preview_fname = tk.StringVar(value="")
        ttk.Label(right, textvariable=self._preview_fname,
                   font=(CN_FONT, 8), foreground='#444').pack(anchor='w', pady=(4, 0))

        # 确定选择按钮
        ttk.Button(right, text="✔ 确定选择", width=15,
                    command=self._confirm_selection).pack(anchor='w', pady=(6, 0))

    def _on_cat_changed(self, event=None):
        titles = TITLES_BY_CATEGORY.get(self.cat_var.get(), ["其他"])
        self.title_cb['values'] = titles
        self.title_var.set(titles[0] if titles else '')

    def _title_right_click(self, event):
        """单张模式标题下拉框右键菜单：删除当前选中标题"""
        title = self.title_var.get().strip()
        if not title:
            return
        menu = tk.Menu(self, tearoff=0)
        menu.add_command(
            label=f"删除标题「{title}」",
            command=lambda: self._do_delete_title(title)
        )
        menu.tk_popup(event.x_root, event.y_root)

    def _do_delete_title(self, title):
        category_ui = self.cat_var.get()
        if not messagebox.askyesno("确认删除",
                f"确定要从词条库中删除标题「{title}」吗？\n（该标题下的所有词条也将一并删除）"):
            return
        self.app.delete_title_from_db(category_ui, title)
        # 刷新下拉框
        new_titles = TITLES_BY_CATEGORY.get(category_ui, [])
        self.title_cb['values'] = new_titles
        self.title_var.set(new_titles[0] if new_titles else '')
        messagebox.showinfo("完成", f"已删除标题「{title}」")

    def _auto_select(self):
        """自动选中当前已选的照片"""
        fname = self.photo_name
        inner = getattr(self, '_pick_inner', None)
        if not fname or not inner:
            return
        for child in inner.winfo_children():
            if getattr(child, '_fname', None) == fname:
                child.event_generate('<Button-1>')
                break

    def _load_thumb_grid(self):
        """加载照片缩略图网格"""
        inner = getattr(self, '_pick_inner', None)
        canvas = getattr(self, '_pick_canvas', None)
        if inner is None or canvas is None:
            return

        for child in inner.winfo_children():
            child.destroy()

        photos = []
        if os.path.exists(PHOTO_FOLDER):
            photos = sorted([
                f for f in os.listdir(PHOTO_FOLDER)
                if f.lower().endswith(('.jpg', '.jpeg', '.png'))
            ])
        
        # 排序：未处理照片排前面，已处理照片排后面
        photos = sorted(photos, key=lambda f: (f in self.processed, f))

        self._status_var.set(f"共 {len(photos)} 张")

        if not photos:
            ttk.Label(inner, text="验收照片文件夹为空",
                       font=(CN_FONT, 10)).grid(row=0, column=0, padx=20, pady=20)
            canvas.configure(scrollregion=canvas.bbox("all"))
            return

        self._photo_list = photos  # 保存照片列表供键盘导航使用
        self._selected_index = -1  # 当前选中索引

        col = 0
        row = 0
        for idx, fname in enumerate(photos):
            frame = ttk.Frame(inner, relief='raised', borderwidth=2,
                              width=self.THUMB_W, height=self.THUMB_H + 36)
            frame.grid(row=row, column=col, padx=6, pady=6)
            frame.pack_propagate(False)
            frame._fname = fname
            frame._index = idx  # 保存索引

            thumb_label = ttk.Label(frame, cursor='hand2')
            thumb_label.pack(fill=tk.X)

            # 异步加载缩略图
            def _load(f=fname, lbl=thumb_label):
                path = os.path.join(PHOTO_FOLDER, f)
                try:
                    img = Image.open(path)
                    img.thumbnail((self.THUMB_W - 10, self.THUMB_H - 10), Image.LANCZOS)
                    thumb = ImageTk.PhotoImage(img)
                    lbl.config(image=thumb, text='')
                    lbl.image = thumb
                except Exception:
                    lbl.config(text='✗', font=(CN_FONT, 16))

            frame.after(10 + idx * 5, _load)

            display_name = fname if len(fname) <= 18 else fname[:15] + '...'
            ttk.Label(frame, text=display_name,
                       font=(CN_FONT, 7), anchor='center').pack(fill=tk.X)

            # 单击：选中预览
            def make_select(fname=fname, fr=frame, idx=idx):
                def select(e=None):
                    for child in inner.winfo_children():
                        child.config(relief='raised', borderwidth=2)
                    fr.config(relief='solid', borderwidth=3)
                    self._show_preview(fname)
                    self._selected_index = idx
                    self._selected_fname_var.set(fname)
                return select

            frame.bind('<Button-1>', make_select(fname, frame, idx))
            thumb_label.bind('<Button-1>', make_select(fname, frame, idx))

            # 双击：直接选择并收起面板（修复闭包问题）
            def make_confirm(fname=fname):
                def confirm(e=None):
                    self._set_photo(fname)
                    self._toggle_thumb_panel()  # 收起面板
                return confirm
            frame.bind('<Double-Button-1>', make_confirm(fname))

            col += 1
            if col >= self.COLS:
                col = 0
                row += 1

        inner.update_idletasks()
        canvas.configure(scrollregion=canvas.bbox("all"))

        # 绑定键盘上下键选择缩略图
        self._bind_keyboard_nav()

    def _bind_keyboard_nav(self):
        """绑定键盘上下键选择缩略图"""
        canvas = getattr(self, '_pick_canvas', None)
        inner = getattr(self, '_pick_inner', None)
        if not canvas or not inner:
            return

        def _on_key_up(e):
            if self._selected_index > 0:
                new_idx = self._selected_index - self.COLS
                if new_idx >= 0:
                    self._select_by_index(new_idx)
            return 'break'  # 阻止默认行为

        def _on_key_down(e):
            if self._selected_index < len(self._photo_list) - 1:
                new_idx = self._selected_index + self.COLS
                if new_idx < len(self._photo_list):
                    self._select_by_index(new_idx)
            return 'break'

        canvas.bind('<Up>', _on_key_up)
        canvas.bind('<Down>', _on_key_down)

        # 让 canvas 获取焦点以接收键盘事件
        canvas.focus_set()

    def _select_by_index(self, idx):
        """根据索引选中缩略图"""
        inner = getattr(self, '_pick_inner', None)
        if not inner or idx < 0 or idx >= len(self._photo_list):
            return

        for child in inner.winfo_children():
            if getattr(child, '_index', -1) == idx:
                child.event_generate('<Button-1>')
                # 确保选中的缩略图可见
                canvas = getattr(self, '_pick_canvas', None)
                if canvas:
                    canvas.see(child.winfo_id())
                break

    def _show_preview(self, fname):
        """在预览区显示照片，并查找匹配的记录填充表单"""
        display = fname[:40] + ('...' if len(fname) > 40 else '')
        self._preview_fname.set(display)
        path = os.path.join(PHOTO_FOLDER, fname)
        try:
            img = Image.open(path)
            img.thumbnail((160, 120), Image.LANCZOS)  # 减小预览图尺寸
            thumb = ImageTk.PhotoImage(img)
            self._preview_label.config(image=thumb, text='')
            self._preview_label.image = thumb
            self._preview_original_path = path
            self.current_photo_name = fname
            # 查找匹配的记录并填充表单
            self.app._lookup_and_fill(fname)
        except Exception:
            self._preview_label.config(text="无法预览\n" + fname, image='')
            self._preview_original_path = None

    def _confirm_selection(self):
        """确认选择当前预览的照片"""
        path = getattr(self, '_preview_original_path', None)
        if not path or not os.path.exists(path):
            messagebox.showinfo("提示", "请先预览一张照片")
            return
        fname = os.path.basename(path)
        self._set_photo(fname)
        self._toggle_thumb_panel()  # 收起面板

    def _open_original_photo(self):
        """双击预览区打开原图"""
        path = getattr(self, '_preview_original_path', None)
        if path and os.path.exists(path):
            _open_file(path)
        else:
            messagebox.showinfo("提示", "请先选择一张照片")

    def _browse_photo(self):
        """浏览并导入照片"""
        from tkinter import filedialog
        paths = filedialog.askopenfilenames(
            parent=self.winfo_toplevel(),
            title="选择照片（可多选）",
            filetypes=[("图片文件", "*.jpg *.jpeg *.png"), ("所有文件", "*.*")]
        )
        if not paths:
            return
        os.makedirs(PHOTO_FOLDER, exist_ok=True)
        copied = []
        for src in paths:
            fname = os.path.basename(src)
            dst = os.path.join(PHOTO_FOLDER, fname)
            if os.path.exists(dst) and not os.path.samefile(src, dst):
                base, ext = os.path.splitext(fname)
                counter = 1
                while os.path.exists(dst):
                    new_name = f"{base}_{counter}{ext}"
                    dst = os.path.join(PHOTO_FOLDER, new_name)
                    counter += 1
                fname = os.path.basename(dst)
            if not os.path.exists(dst) or not os.path.samefile(src, dst):
                import shutil
                shutil.copy2(src, dst)
            copied.append(fname)

        self._load_thumb_grid()
        messagebox.showinfo("完成", f"已添加 {len(copied)} 张照片到验收照片文件夹")

    def _set_photo(self, fname):
        """设置选中的照片"""
        self.photo_name = fname
        self.photo_path = os.path.join(PHOTO_FOLDER, fname)
        try:
            img = Image.open(self.photo_path)
            img.thumbnail((120, 90))
            self._thumb = ImageTk.PhotoImage(img)
            self.thumb_label.config(image=self._thumb, text='', background='white')
        except Exception:
            self.thumb_label.config(text=fname, image='', background='#e0e0e0')

    def _open_selected_photo(self):
        """双击已选照片标签打开原图"""
        if self.photo_name:
            path = os.path.join(PHOTO_FOLDER, self.photo_name)
            if os.path.exists(path):
                _open_file(path)

    def _pick_desc(self):
        self.app.open_db_picker(
            parent=self.winfo_toplevel(),
            title="选择问题描述",
            db=self.app.issues_db,
            category=self.cat_var.get(),
            issue_title=self.title_var.get(),
            callback=lambda v: (self.desc_text.delete('1.0', tk.END),
                                self.desc_text.insert('1.0', v))
        )

    def _pick_remedy(self):
        self.app.open_db_picker(
            parent=self.winfo_toplevel(),
            title="选择整改意见",
            db=self.app.remedies_db,
            category=self.cat_var.get(),
            issue_title=self.title_var.get(),
            callback=lambda v: (self.remedy_text.delete('1.0', tk.END),
                                self.remedy_text.insert('1.0', v))
        )

    def get_record(self):
        return {
            'photo': self.photo_name,
            'category': self.cat_var.get(),
            'title': self.title_var.get(),
            'severity': self.sev_var.get(),
            'description': self.desc_text.get('1.0', tk.END).strip(),
            'remedy': self.remedy_text.get('1.0', tk.END).strip(),
            'plan_img': self.plan_var.get(),
        }

    def clear(self):
        self.photo_path = None
        self.photo_name = None
        self._thumb = None
        self.thumb_label.config(image='', text='未选照片', background='#e0e0e0')
        self.cat_var.set(CATEGORIES[0])
        titles = TITLES_BY_CATEGORY.get(CATEGORIES[0], [])
        self.title_cb['values'] = titles
        self.title_var.set(titles[0] if titles else '')
        self.sev_var.set(SEVERITIES[0])
        self.desc_text.delete('1.0', tk.END)
        self.remedy_text.delete('1.0', tk.END)
        self.plan_var.set('无')


class PhotoReportApp:
    def __init__(self, root):
        self.root = root
        self.root.title("华为门店验收工具 v2.8")
        self.root.geometry("1200x900")
        self.root.resizable(True, True)

        # ====== 菜单栏 ======
        self._build_menu()

        # 数据
        self.photos = []
        self.current_idx = 0
        self.session_data = load_json(SESSION_FILE, [])
        self.processed = {item['photo'] for item in self.session_data}

        self.issues_db = load_db(ISSUES_DB_FILE)
        self.remedies_db = load_db(REMEDIES_DB_FILE)
        self.plan_images = self.get_plan_images()

        self.settings = load_settings()

        self.load_photos()
        self.current_record = {}
        self.recorded_data = []  # 存储已记录的幻灯片数据
        self._thumb_frames = {}  # fname -> frame 映射，用于刷新缩略图

        self.create_widgets()
        self.show_photo(0)

    # ====== 菜单栏 ======
    def _build_menu(self):
        """构建顶部菜单栏"""
        menubar = tk.Menu(self.root)
        self.root.config(menu=menubar)

        # 帮助菜单
        help_menu = tk.Menu(menubar, tearoff=0)
        menubar.add_cascade(label="帮助", menu=help_menu)
        help_menu.add_command(label="关于", command=self._show_about)

    def _show_about(self):
        """显示关于对话框"""
        about = tk.Toplevel(self.root)
        about.title("关于")
        about.geometry("400x220")
        about.resizable(False, False)
        about.transient(self.root)
        about.grab_set()

        # 居中显示
        about.update_idletasks()
        x = self.root.winfo_x() + (self.root.winfo_width() // 2) - 200
        y = self.root.winfo_y() + (self.root.winfo_height() // 2) - 110
        about.geometry(f"+{x}+{y}")

        tk.Label(about, text="华为门店验收工具", font=(CN_FONT, 16, 'bold')).pack(pady=(20, 5))
        tk.Label(about, text="版本：v2.8", font=(CN_FONT, 10)).pack()
        tk.Label(about, text="作者：王柱勇", font=(CN_FONT, 10)).pack(pady=(10, 2))
        tk.Label(about, text="联系电话：13945079480", font=(CN_FONT, 10)).pack()
        tk.Label(about, text="有更好的改进建议，请联系作者", font=(CN_FONT, 9), fg="gray").pack(pady=(10, 5))

        ttk.Button(about, text="确定", command=about.destroy, width=12).pack(pady=(5, 15))

    # ====== 通用工具 ======

    def get_plan_images(self):
        images = ["无"]
        if os.path.exists(PLAN_IMAGES_FOLDER):
            for f in os.listdir(PLAN_IMAGES_FOLDER):
                if f.lower().endswith(('.png', '.jpg', '.jpeg')):
                    images.append(f)
        return images

    def load_photos(self):
        if os.path.exists(PHOTO_FOLDER):
            photos = sorted([
                f for f in os.listdir(PHOTO_FOLDER)
                if f.lower().endswith(('.jpg', '.jpeg', '.png'))
            ])
            # 排序：未处理照片排前面，已处理照片排后面
            self.photos = sorted(photos, key=lambda f: (f in self.processed, f))
        if not self.photos:
            messagebox.showwarning("提示", f"照片文件夹为空：\n{PHOTO_FOLDER}")

    def _browse_photo(self):
        """浏览并导入照片"""
        from tkinter import filedialog
        paths = filedialog.askopenfilenames(
            parent=self.root,
            title="选择照片（可多选）",
            filetypes=[("图片文件", "*.jpg *.jpeg *.png"), ("所有文件", "*.*")]
        )
        if not paths:
            return
        os.makedirs(PHOTO_FOLDER, exist_ok=True)
        copied = []
        for src in paths:
            fname = os.path.basename(src)
            dst = os.path.join(PHOTO_FOLDER, fname)
            if os.path.exists(dst) and not os.path.samefile(src, dst):
                base, ext = os.path.splitext(fname)
                counter = 1
                while os.path.exists(dst):
                    new_name = f"{base}_{counter}{ext}"
                    dst = os.path.join(PHOTO_FOLDER, new_name)
                    counter += 1
                fname = os.path.basename(dst)
            if not os.path.exists(dst) or not os.path.samefile(src, dst):
                import shutil
                shutil.copy2(src, dst)
            copied.append(fname)
        self.load_photos()
        self._load_single_thumb_grid()
        messagebox.showinfo("完成", f"已添加 {len(copied)} 张照片到验收照片文件夹")

    def _import_app_json(self):
        """从App导入JSON数据，自动匹配照片并填充记录"""
        from tkinter import filedialog
        path = filedialog.askopenfilename(
            parent=self.root,
            title="选择App导出的JSON文件",
            filetypes=[("JSON文件", "*.json"), ("所有文件", "*.*")]
        )
        if not path:
            return

        try:
            with open(path, 'r', encoding='utf-8') as f:
                app_data = json.load(f)
        except Exception as e:
            messagebox.showerror("导入失败", f"无法读取JSON文件：\n{e}")
            return

        if not isinstance(app_data, list):
            messagebox.showerror("格式错误", "JSON文件格式不正确，应为记录列表")
            return

        matched = 0
        skipped = 0
        debug_info = []
        
        # 建立PC端照片文件名映射
        pc_photos_by_name = {}
        if os.path.exists(PHOTO_FOLDER):
            for f in os.listdir(PHOTO_FOLDER):
                if f.lower().endswith(('.jpg', '.jpeg', '.png', '.webp')):
                    pc_photos_by_name[f.lower()] = f
        
        def find_best_match(app_photo_path, record_idx):
            """根据文件名匹配PC端照片"""
            basename = os.path.basename(app_photo_path)
            if basename.lower() in pc_photos_by_name:
                return pc_photos_by_name[basename.lower()]
            debug_info.append(f"  [{record_idx}] ✗ 未找到: {basename}")
            return None
        
        for idx, record in enumerate(app_data):
            # 兼容App端字段名（photo / photoPath / fileName）
            photo_name = record.get('photo', '') or record.get('photoPath', '') or record.get('fileName', '')
            if not photo_name:
                skipped += 1
                debug_info.append(f"[{idx}] [跳过] 记录无photo/fileName字段: {str(record)[:60]}")
                continue

            # 根据时间匹配PC端照片
            matched_name = find_best_match(photo_name, idx)
            
            if not matched_name:
                debug_info.append(f"[未找到] APP: '{os.path.basename(photo_name)}'")
                skipped += 1
                continue

            debug_info.append(f"[按文件名匹配] '{os.path.basename(photo_name)}' -> '{matched_name}'")
            std_record = {
                'photo': matched_name,
                'category': record.get('category', '店铺问题—基装'),
                'title': record.get('title', ''),
                'severity': record.get('severity', '一般问题'),
                'description': record.get('description', ''),
                'remedy': record.get('remedy', ''),
                'plan_img': record.get('plan_img', '') or record.get('planImage', '无'),
            }

            self.recorded_data = [r for r in self.recorded_data if r['photo'] != matched_name]
            self.recorded_data.append(std_record)
            self.processed.add(matched_name)
            self._refresh_thumbnail(matched_name)
            matched += 1

        self.record_count_label.config(text=f"已记录: {len(self.recorded_data)} 张")
        self.stats_label.config(text=f"已处理: {len(self.processed)}/{len(self.photos)} 张")

        msg = f"导入完成\n\n匹配成功: {matched} 条\n跳过(无照片): {skipped} 条"

        # 显示调试信息
        if debug_info:
            debug_msg = "【调试信息】\n\n" + "\n".join(debug_info[:20])  # 最多显示20条
            if len(debug_info) > 20:
                debug_msg += f"\n... 还有 {len(debug_info)-20} 条"
            messagebox.showinfo("导入详情", debug_msg)

        if matched > 0:
            messagebox.showinfo("导入成功", msg)
        else:
            messagebox.showwarning("导入结果", msg + "\n\n提示：请确认照片已放入「验收照片」文件夹")

    def _load_single_thumb_grid(self):
        """加载单张录入标签页的照片缩略图"""
        inner = getattr(self, '_single_pick_inner', None)
        canvas = getattr(self, '_single_pick_canvas', None)
        if inner is None or canvas is None:
            return

        for child in inner.winfo_children():
            child.destroy()

        photos = []
        if os.path.exists(PHOTO_FOLDER):
            photos = sorted([
                f for f in os.listdir(PHOTO_FOLDER)
                if f.lower().endswith(('.jpg', '.jpeg', '.png'))
            ])
        photos = sorted(photos, key=lambda f: (f in self.processed, f))

        if not photos:
            ttk.Label(inner, text="验收照片文件夹为空",
                       font=(CN_FONT, 10)).grid(row=0, column=0, padx=20, pady=20)
            canvas.configure(scrollregion=canvas.bbox("all"))
            return

        COLS = 4
        THUMB_W = 180
        THUMB_H = 140

        self._thumb_frames.clear()
        self._single_photo_list = photos
        self._single_selected_index = -1

        for idx, fname in enumerate(photos):
            r, c = divmod(idx, COLS)
            frame = tk.Frame(inner, bd=0, bg='#E0E0E0', cursor='hand2',
                              highlightbackground='#E0E0E0', highlightthickness=3)
            frame.grid(row=r, column=c, padx=5, pady=5, sticky='nsew')
            inner.columnconfigure(c, weight=1)

            img_path = os.path.join(PHOTO_FOLDER, fname)
            photo = None
            try:
                img = Image.open(img_path)
                img.thumbnail((THUMB_W, THUMB_H), Image.LANCZOS)
                draw = ImageDraw.Draw(img)
                if fname in self.processed:
                    w, h = img.size
                    r_size = 15
                    x0, y0 = w - r_size*2 - 5, 5
                    x1, y1 = x0 + r_size*2, y0 + r_size*2
                    draw.ellipse([x0, y0, x1, y1], fill='#00AA00', outline='white', width=2)
                    cx, cy = (x0 + x1)//2, (y0 + y1)//2
                    draw.line([cx-5, cy, cx, cy+5, cx+8, cy-8], fill='white', width=3)
                photo = ImageTk.PhotoImage(img)
            except Exception:
                pass

            thumb_label = tk.Label(frame, image=photo, text='',
                                   compound='top', anchor='center',
                                   font=(CN_FONT, 7), wraplength=THUMB_W,
                                   bg='#F5F5F5')
            if photo:
                thumb_label.image = photo
            thumb_label.pack(fill=tk.BOTH, expand=True)

            self._thumb_frames[fname] = frame
            frame._thumb_label = thumb_label

            def _select(e=None, ff=fname, fr=frame, idx=idx):
                for fn2, frm2 in self._thumb_frames.items():
                    frm2.configure(highlightbackground='#AAAAAA', bg='#E0E0E0')
                fr.configure(highlightbackground='#FF0000', bg='#FFCCCC')
                self.current_photo_name = ff
                self._single_selected_index = idx
                self._select_single_photo(ff)
                self._lookup_and_fill(ff)

            def _open(e=None, ff=fname):
                path = os.path.join(PHOTO_FOLDER, ff)
                if os.path.exists(path):
                    _open_file(path)

            thumb_label.bind('<Button-1>', _select)
            thumb_label.bind('<Double-Button-1>', _open)

        # 更新计数器
        self.stats_label.config(text=f"已处理: {len(self.processed)}/{len(self.photos)} 张")
        canvas.configure(scrollregion=canvas.bbox("all"))

    def _lookup_and_fill(self, fname):
        """点选照片时，根据文件名查找 recorded_data 中的记录并填充表单；
        无匹配记录时自动复用上一条记录的所有字段（仅换照片）"""
        for r in getattr(self, 'recorded_data', []):
            if r.get('photo') == fname:
                self._fill_form_with_record(r)
                return
        # 没找到匹配记录 → 自动复用上一条记录（快捷模式）
        if self.recorded_data:
            last = dict(self.recorded_data[-1])  # 浅拷贝
            last['photo'] = fname  # 只换照片
            self._fill_form_with_record(last)
        else:
            self._clear_form()

    def _clear_form(self):
        """清空左侧表单"""
        self.category_var.set(CATEGORIES[0])
        self._on_category_changed()
        self.title_var.set(TITLES_BY_CATEGORY[CATEGORIES[0]][0])
        self.severity_var.set(SEVERITIES[0])
        self.description_text.delete('1.0', tk.END)
        self.remedy_text.delete('1.0', tk.END)
        self.plan_var.set('无')

    def _fill_form_with_record(self, record):
        """将记录数据填充到左侧表单并选中对应照片"""
        pass
        try:
            # 填充表单
            self.category_var.set(record.get('category', CATEGORIES[0]))
            self._on_category_changed()

            # 安全设置标题：如果标题不在当前分类列表中则跳过
            target_title = record.get('title', '')
            current_titles = TITLES_BY_CATEGORY.get(self.category_var.get(), [])
            if target_title in current_titles:
                self.title_var.set(target_title)
            elif current_titles:
                self.title_var.set(current_titles[0])

            self.severity_var.set(record.get('severity', SEVERITIES[0]))
            self.description_text.delete('1.0', tk.END)
            self.description_text.insert('1.0', record.get('description', ''))
            self.remedy_text.delete('1.0', tk.END)
            self.remedy_text.insert('1.0', record.get('remedy', ''))
            # 兼容 planImage / plan_img 字段，值不带后缀时尝试补 .png
            plan = record.get('plan_img', '') or record.get('planImage', '无')
            if plan and plan != '无':
                # 尝试直接匹配，或补 .png 后缀匹配
                if plan in self.plan_images:
                    self.plan_var.set(plan)
                else:
                    found = False
                    for pi in self.plan_images:
                        if pi == plan or pi == plan + '.png' or pi.rstrip('.png') == plan:
                            self.plan_var.set(pi)
                            found = True
                            break
                    if not found:
                        self.plan_var.set('无')
            else:
                self.plan_var.set('无')

            # 选中对应照片（在缩略图中高亮）
            photo_name = record.get('photo', '')
            if photo_name:
                self.current_photo_name = photo_name
                if hasattr(self, '_thumb_frames') and photo_name in self._thumb_frames:
                    frame = self._thumb_frames[photo_name]
                    for fn2, frm2 in self._thumb_frames.items():
                        frm2.configure(highlightbackground='#AAAAAA', bg='#E0E0E0')
                    frame.configure(highlightbackground='#FF0000', bg='#FFCCCC')
                self.photo_name_label.config(text=f"已选择：{photo_name}", foreground='#000')
            pass
        except Exception as e:
            pass
            import traceback
            traceback.print_exc()

    def _processed_style(self, fname):
        """根据照片处理状态返回背景色（已废弃，仅兼容保留）"""
        return '#CCEECC' if fname in self.processed else '#CCCCCC'

    def _select_single_photo(self, fname):
        """选中单张录入标签页的照片"""
        self.current_photo_name = fname
        self.photo_name_label.config(text=fname)
        # 缩略图已显示选中状态，无需额外预览区

    def open_db_picker(self, parent, title, db, category, issue_title, callback):
        """通用数据库选择弹窗（支持编辑/删除条目）"""
        cat_key = category.replace("店铺问题—", "").replace("店铺问题-", "")

        # 判断是问题库还是整改库
        is_issues = (db is self.issues_db)
        db_file = ISSUES_DB_FILE if is_issues else REMEDIES_DB_FILE

        def _get_items():
            items = []
            if cat_key in db:
                items = db[cat_key].get(issue_title, [])
            if not items:
                items = db.get(issue_title, [])
            return list(items)

        if not _get_items():
            messagebox.showinfo("提示", f"「{issue_title}」暂无库词条，请手动输入")
            return

        win = tk.Toplevel(parent)
        win.title(title)
        win.geometry("560x400")
        win.grab_set()

        # 说明
        ttk.Label(win, text=f"分类：{cat_key}  /  标题：{issue_title}",
                  font=(CN_FONT, 9), foreground='#555').pack(anchor='w', padx=10, pady=(6, 0))

        list_frame = ttk.Frame(win)
        list_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=6)

        vsb = ttk.Scrollbar(list_frame, orient=tk.VERTICAL)
        vsb.pack(side=tk.RIGHT, fill=tk.Y)
        lb = tk.Listbox(list_frame, selectmode=tk.SINGLE, font=(CN_FONT, 10),
                        yscrollcommand=vsb.set, activestyle='dotbox')
        lb.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        vsb.config(command=lb.yview)

        def _refresh_list():
            items = _get_items()
            lb.delete(0, tk.END)
            for it in items:
                lb.insert(tk.END, it)

        _refresh_list()

        def _save_db():
            save_json(db_file, db)

        def on_ok(e=None):
            sel = lb.curselection()
            if sel:
                callback(_get_items()[sel[0]])
            win.destroy()

        def on_edit():
            sel = lb.curselection()
            if not sel:
                messagebox.showinfo("提示", "请先选中一条词条", parent=win)
                return
            idx = sel[0]
            items = _get_items()
            old_val = items[idx]

            edit_win = tk.Toplevel(win)
            edit_win.title("编辑词条")
            edit_win.geometry("500x200")
            edit_win.grab_set()

            ttk.Label(edit_win, text="修改内容：", font=(CN_FONT, 9)).pack(anchor='w', padx=10, pady=(10, 0))
            txt = tk.Text(edit_win, width=55, height=5, font=(CN_FONT, 10))
            txt.pack(padx=10, pady=5)
            txt.insert('1.0', old_val)
            txt.focus_set()

            def on_save():
                new_val = txt.get('1.0', tk.END).strip()
                if not new_val:
                    messagebox.showwarning("提示", "内容不能为空", parent=edit_win)
                    return
                # 更新数据库
                if cat_key in db and issue_title in db[cat_key]:
                    target = db[cat_key][issue_title]
                elif issue_title in db:
                    target = db[issue_title]
                else:
                    target = None
                if target is not None and old_val in target:
                    pos = target.index(old_val)
                    target[pos] = new_val
                    _save_db()
                edit_win.destroy()
                _refresh_list()

            btn_row = ttk.Frame(edit_win)
            btn_row.pack(pady=4)
            ttk.Button(btn_row, text="保存", command=on_save).pack(side=tk.LEFT, padx=6)
            ttk.Button(btn_row, text="取消", command=edit_win.destroy).pack(side=tk.LEFT, padx=6)

        def on_delete():
            sel = lb.curselection()
            if not sel:
                messagebox.showinfo("提示", "请先选中一条词条", parent=win)
                return
            idx = sel[0]
            items = _get_items()
            val = items[idx]
            if not messagebox.askyesno("确认删除", f"确定删除以下词条？\n\n{val}", parent=win):
                return
            # 从数据库删除
            if cat_key in db and issue_title in db[cat_key]:
                target = db[cat_key][issue_title]
            elif issue_title in db:
                target = db[issue_title]
            else:
                target = None
            if target is not None and val in target:
                target.remove(val)
                _save_db()
            _refresh_list()

        # 底部按钮
        btn_frame = ttk.Frame(win)
        btn_frame.pack(fill=tk.X, padx=10, pady=(0, 8))

        ttk.Button(btn_frame, text="✔ 选择使用", command=on_ok).pack(side=tk.LEFT, padx=(0, 8))
        ttk.Button(btn_frame, text="✎ 编辑", command=on_edit).pack(side=tk.LEFT, padx=(0, 4))
        ttk.Button(btn_frame, text="✕ 删除", command=on_delete).pack(side=tk.LEFT)

        lb.bind('<Double-Button-1>', on_ok)

    # ====== 单张模式 ======

    def _on_category_changed(self, event=None):
        category = self.category_var.get()
        titles = TITLES_BY_CATEGORY.get(category, ["其他"])
        self.title_combo['values'] = titles
        if titles:
            self.title_var.set(titles[0])
        self._save_selection()

    def _title_right_click_batch(self, event):
        """批量模式标题下拉框右键菜单：删除当前选中标题"""
        title = self.title_var.get().strip()
        if not title:
            return
        menu = tk.Menu(self.root, tearoff=0)
        menu.add_command(
            label=f"删除标题「{title}」",
            command=lambda: self._do_delete_title_batch(title)
        )
        menu.tk_popup(event.x_root, event.y_root)

    def _do_delete_title_batch(self, title):
        category_ui = self.category_var.get()
        if not messagebox.askyesno("确认删除",
                f"确定要从词条库中删除标题「{title}」吗？\n（该标题下的所有词条也将一并删除）",
                parent=self.root):
            return
        self.delete_title_from_db(category_ui, title)
        new_titles = TITLES_BY_CATEGORY.get(category_ui, [])
        self.title_combo['values'] = new_titles
        self.title_var.set(new_titles[0] if new_titles else '')
        messagebox.showinfo("完成", f"已删除标题「{title}」", parent=self.root)

    def _on_title_changed(self, event=None):
        self._save_selection()

    def _on_severity_changed(self, event=None):
        self._save_selection()

    def _save_selection(self):
        self.settings['category'] = self.category_var.get()
        self.settings['title'] = self.title_var.get()
        self.settings['severity'] = self.severity_var.get()
        save_settings(self.settings)

    def show_photo(self, idx):
        if 0 <= idx < len(self.photos):
            self.current_idx = idx
            filename = self.photos[idx]

            # 更新标题
            status = "[已处理]" if filename in self.processed else "[待处理]"
            self.title_label.config(text=f"{status} 第 {idx+1}/{len(self.photos)} 张：{filename}")

            # 更新选中照片名称显示
            if hasattr(self, 'photo_name_label'):
                self.photo_name_label.config(text=f"已选择：{filename}", foreground='#000')

            self.load_record(filename)

    def load_record(self, filename):
        saved_record = None
        for record in self.session_data:
            if record['photo'] == filename:
                saved_record = record.copy()
                break

        if saved_record:
            self.current_record = saved_record
            self.category_var.set(saved_record.get('category', CATEGORIES[0]))
            self.title_var.set(saved_record.get('title', ''))
            self.severity_var.set(saved_record.get('severity', SEVERITIES[0]))
            titles = TITLES_BY_CATEGORY.get(self.category_var.get(), ["其他"])
            self.title_combo['values'] = titles
            self.description_text.delete('1.0', tk.END)
            self.description_text.insert('1.0', saved_record.get('description', ''))
            self.remedy_text.delete('1.0', tk.END)
            self.remedy_text.insert('1.0', saved_record.get('remedy', ''))
            plan = saved_record.get('plan_img', '无')
            if plan in self.plan_images:
                self.plan_var.set(plan)
            else:
                self.plan_var.set('无')
        else:
            self.current_record = {'photo': filename}

    def _refresh_thumbnail(self, fname):
        """刷新指定照片的缩略图（重新加载并绘制已处理标记）"""
        frame = self._thumb_frames.get(fname)
        if not frame:
            return
        thumb_label = getattr(frame, '_thumb_label', None)
        if not thumb_label:
            return
        try:
            from PIL import ImageDraw
            img = Image.open(os.path.join(PHOTO_FOLDER, fname))
            img.thumbnail((200, 150), Image.LANCZOS)
            if fname in self.processed:
                draw = ImageDraw.Draw(img)
                w, h = img.size
                r = 15
                x0, y0 = w - r*2 - 5, 5
                x1, y1 = x0 + r*2, y0 + r*2
                draw.ellipse([x0, y0, x1, y1], fill='#00AA00', outline='white', width=2)
                cx, cy = (x0 + x1) // 2, (y0 + y1) // 2
                draw.line([cx-5, cy, cx, cy+5, cx+8, cy-8], fill='white', width=3)
            thumb = ImageTk.PhotoImage(img)
            thumb_label.config(image=thumb, text='')
            thumb_label.image = thumb
        except Exception:
            pass

    def _refresh_rec_listbox(self):
        """刷新已记录列表控件"""
        lb = getattr(self, '_rec_listbox', None)
        if lb is None:
            return
        lb.delete(0, tk.END)
        for i, rec in enumerate(self.recorded_data):
            photo = rec.get('photo', '（无照片）')
            title = rec.get('title', '')
            severity = rec.get('severity', '')
            lb.insert(tk.END, f"{i+1}. {photo}  [{title} · {severity}]")
        # 滚到最新一条
        if self.recorded_data:
            lb.see(tk.END)

    def record_current(self):
        """记录当前表单数据（不生成PPT）"""
        if not hasattr(self, '_single_selected_index') or self._single_selected_index < 0:
            messagebox.showwarning("提示", "请先选择一张照片")
            return

        record = {
            'photo': self._single_photo_list[self._single_selected_index],
            'category': self.category_var.get(),
            'title': self.title_var.get(),
            'severity': self.severity_var.get(),
            'description': self.description_text.get('1.0', tk.END).strip(),
            'remedy': self.remedy_text.get('1.0', tk.END).strip(),
            'plan_img': self.plan_var.get()
        }

        # 去重：同一张照片只保留最新记录
        self.recorded_data = [r for r in self.recorded_data if r['photo'] != record['photo']]
        self.recorded_data.append(record)

        # 自动写入数据库
        self._auto_save_to_db(record)
        self.processed.add(record['photo'])
        # 刷新缩略图绿勾
        self._refresh_thumbnail(record['photo'])

        # 更新计数
        self.record_count_label.config(text=f"已记录: {len(self.recorded_data)} 张")
        self.stats_label.config(text=f"已处理: {len(self.processed)}/{len(self.photos)} 张")
        # 刷新已记录列表
        self._refresh_rec_listbox()

        self.status_label.config(text=f"[已记录] {record['photo']}", fg="#FF6600")
        self.root.after(2000, lambda: self.status_label.config(text="就绪"))

    def save_current(self):
        """根据 recorded_data 生成多页PPT"""
        if not self.recorded_data:
            messagebox.showwarning("提示", "请先点击'记录'按钮记录至少一张照片")
            return

        records = list(self.recorded_data)
        try:
            os.makedirs(OUTPUT_FOLDER, exist_ok=True)

            prs = Presentation(TEMPLATE_PPTX)
            # 第一页
            self._update_slide_content(prs.slides[0], records[0])
            self._add_decoration(prs.slides[0])  # 嵌入装饰图到幻灯片本身
            self._auto_save_to_db(records[0])

            # 追加后续页
            for rec in records[1:]:
                self._auto_save_to_db(rec)
                tmp_prs = Presentation(TEMPLATE_PPTX)
                tmp_slide = tmp_prs.slides[0]
                self._update_slide_content(tmp_slide, rec)

                new_slide = prs.slides.add_slide(
                    prs.slides[0].slide_layout  # 用第1页同布局(=Master2)，继承MAKE it POSSIBLE装饰
                )
                # 清空布局占位符（不清背景，保留母版的MAKE it POSSIBLE等装饰）
                for shape in list(new_slide.shapes):
                    shape._element.getparent().remove(shape._element)
                # 复制 tmp_slide 的形状（模板形状自带白底，会覆盖母版内容区）
                for shape in tmp_slide.shapes:
                    self._copy_shape_to_slide(shape, tmp_slide, new_slide, prs)
                self._add_decoration(new_slide)  # 嵌入装饰图到幻灯片本身

            now = datetime.now().strftime('%Y%m%d_%H%M%S')
            out_name = f"验收报告_{now}.pptx"
            out_path = os.path.join(OUTPUT_FOLDER, out_name)
            prs.save(out_path)

            _open_file(out_path)

            count = len(records)
            self.record_count_label.config(text="已记录: 0 张")
            self.status_label.config(text=f"[生成完成] {out_name}  共{count}页", fg="green")
            self.root.after(4000, lambda: self.status_label.config(text="就绪"))

            # 清空记录
            self.recorded_data.clear()
            self._refresh_rec_listbox()
            # 重新加载缩略图（已处理的排后面）
            self._load_single_thumb_grid()

        except Exception as e:
            messagebox.showerror("生成失败", f"无法生成PPT文件：\n{str(e)}")

    def _auto_save_to_db(self, record):
        """自动将新词条写入数据库（去重）"""
        category = record['category'].replace("店铺问题—", "").replace("店铺问题-", "")
        title = record['title'].strip()

        # 标题不在预设则追加
        if title and title not in TITLES_BY_CATEGORY.get(record['category'], []):
            TITLES_BY_CATEGORY[record['category']].append(title)
            try:
                self.title_combo['values'] = TITLES_BY_CATEGORY[record['category']]
            except Exception:
                pass

        desc = record['description']
        if desc:
            if category not in self.issues_db:
                self.issues_db[category] = {}
            if title not in self.issues_db[category]:
                self.issues_db[category][title] = []
            if desc not in self.issues_db[category][title]:
                self.issues_db[category][title].append(desc)
                save_json(ISSUES_DB_FILE, self.issues_db)

        remedy = record['remedy']
        if remedy:
            if category not in self.remedies_db:
                self.remedies_db[category] = {}
            if title not in self.remedies_db[category]:
                self.remedies_db[category][title] = []
            if remedy not in self.remedies_db[category][title]:
                self.remedies_db[category][title].append(remedy)
                save_json(REMEDIES_DB_FILE, self.remedies_db)

    def generate_single_slide(self, record):
        try:
            os.makedirs(OUTPUT_FOLDER, exist_ok=True)

            prs = Presentation(TEMPLATE_PPTX)
            slide = prs.slides[0]

            self._update_slide_content(slide, record)
            self._add_decoration(slide)  # 嵌入装饰图到幻灯片本身

            title = record.get('title', '未分类').strip()
            safe_title = title.replace('/', '-').replace('\\', '-').replace(':', '-')
            filename = f"幻灯片_{safe_title}.pptx" if safe_title else f"幻灯片_{record['photo']}"

            output_path = os.path.join(OUTPUT_FOLDER, filename)
            if os.path.exists(output_path):
                base, ext = os.path.splitext(filename)
                counter = 1
                while os.path.exists(output_path):
                    filename = f"{base}_{counter}.pptx"
                    output_path = os.path.join(OUTPUT_FOLDER, filename)
                    counter += 1

            prs.save(output_path)

            _open_file(output_path)

            return output_path
        except Exception as e:
            messagebox.showerror("生成失败", f"无法生成PPT文件：\n{str(e)}")
            return None

    def prev_photo(self):
        if self.current_idx > 0:
            self.show_photo(self.current_idx - 1)

    def next_photo(self):
        if self.current_idx < len(self.photos) - 1:
            self.show_photo(self.current_idx + 1)

    def save_and_next(self):
        self.save_current()
        if self.current_idx < len(self.photos) - 1:
            self.show_photo(self.current_idx + 1)

    def add_to_issues(self):
        desc = self.description_text.get('1.0', tk.END).strip()
        if not desc:
            messagebox.showwarning("提示", "请先填写问题描述内容")
            return
        category = self.category_var.get().replace("店铺问题—", "").replace("店铺问题-", "")
        title = self.title_var.get().strip()
        if not title:
            messagebox.showwarning("提示", "请先选择标题")
            return
        if category not in self.issues_db:
            self.issues_db[category] = {}
        if title not in self.issues_db[category]:
            self.issues_db[category][title] = []
        if desc not in self.issues_db[category][title]:
            self.issues_db[category][title].append(desc)
            save_json(ISSUES_DB_FILE, self.issues_db)
            messagebox.showinfo("提示", f"已添加到问题库\n[{category} > {title}]：\n{desc}")
        else:
            messagebox.showinfo("提示", "该词条已存在于问题库中")

    def add_to_remedies(self):
        remedy = self.remedy_text.get('1.0', tk.END).strip()
        if not remedy:
            messagebox.showwarning("提示", "请先填写整改意见内容")
            return
        category = self.category_var.get().replace("店铺问题—", "").replace("店铺问题-", "")
        title = self.title_var.get().strip()
        if not title:
            messagebox.showwarning("提示", "请先选择标题")
            return
        if category not in self.remedies_db:
            self.remedies_db[category] = {}
        if title not in self.remedies_db[category]:
            self.remedies_db[category][title] = []
        if remedy not in self.remedies_db[category][title]:
            self.remedies_db[category][title].append(remedy)
            save_json(REMEDIES_DB_FILE, self.remedies_db)
            messagebox.showinfo("提示", f"已添加到整改库\n[{category} > {title}]：\n{remedy}")
        else:
            messagebox.showinfo("提示", "该词条已存在于整改库中")

    def select_from_issues(self):
        self.open_db_picker(
            parent=self.root,
            title="选择问题描述",
            db=self.issues_db,
            category=self.category_var.get(),
            issue_title=self.title_var.get(),
            callback=lambda v: (self.description_text.delete('1.0', tk.END),
                                self.description_text.insert('1.0', v))
        )

    def select_from_remedies(self):
        self.open_db_picker(
            parent=self.root,
            title="选择整改意见",
            db=self.remedies_db,
            category=self.category_var.get(),
            issue_title=self.title_var.get(),
            callback=lambda v: (self.remedy_text.delete('1.0', tk.END),
                                self.remedy_text.insert('1.0', v))
        )

    def delete_title_from_db(self, category_ui, title):
        """从 issues_db / remedies_db 中删除指定标题（含其所有词条），并刷新全局 TITLES_BY_CATEGORY"""
        if not title:
            return
        cat_key = category_ui.replace("店铺问题—", "").replace("店铺问题-", "")
        changed = False
        for db, db_file in [(self.issues_db, ISSUES_DB_FILE),
                            (self.remedies_db, REMEDIES_DB_FILE)]:
            if cat_key in db and title in db[cat_key]:
                del db[cat_key][title]
                save_json(db_file, db)
                changed = True
        if changed:
            # 刷新全局标题列表
            TITLES_BY_CATEGORY.update(_build_titles_from_db())

    # ====== 批量模式 ======

    def _build_batch_slots(self, count):
        """根据 count 重建批量槽位"""
        for slot in self.batch_slots:
            slot.destroy()
        self.batch_slots = []

        for i in range(count):
            slot = BatchSlot(self.batch_inner_frame, index=i, app=self)
            slot.pack(fill=tk.X, pady=4, padx=6)
            self.batch_slots.append(slot)

        # 更新滚动区域
        self.batch_inner_frame.update_idletasks()
        self.batch_canvas.configure(scrollregion=self.batch_canvas.bbox("all"))

    def _on_batch_count_apply(self):
        try:
            n = int(self.batch_count_var.get())
            if n < 1 or n > 50:
                raise ValueError
        except ValueError:
            messagebox.showwarning("提示", "请输入 1~50 之间的数字")
            return
        self._build_batch_slots(n)

    def batch_clear_all(self):
        if messagebox.askyesno("确认", "清空所有槽位内容？"):
            for slot in self.batch_slots:
                slot.clear()

    def batch_generate(self):
        """收集所有有照片的槽位，一次性生成多页 PPT"""
        records = []
        for slot in self.batch_slots:
            r = slot.get_record()
            if r['photo']:   # 只处理已选照片的槽位
                records.append(r)

        if not records:
            messagebox.showwarning("提示", "请至少为一个槽位选择照片")
            return

        if not messagebox.askyesno("确认生成",
                                   f"共 {len(records)} 条记录，生成一个 {len(records)} 页 PPT？"):
            return

        try:
            os.makedirs(OUTPUT_FOLDER, exist_ok=True)

            import shutil
            # 以第一条生成单页 PPT 作为底稿
            prs = Presentation(TEMPLATE_PPTX)
            first_slide = prs.slides[0]
            self._update_slide_content(first_slide, records[0])
            self._add_decoration(first_slide)  # 嵌入装饰图到幻灯片本身

            # 自动入库
            self._auto_save_to_db(records[0])

            # 追加后续页
            for rec in records[1:]:
                self._auto_save_to_db(rec)
                # 新建一个幻灯片（从模板另存临时文件获取干净 slide）
                tmp_prs = Presentation(TEMPLATE_PPTX)
                tmp_slide = tmp_prs.slides[0]
                self._update_slide_content(tmp_slide, rec)

                # 把 tmp_slide 追加到主 prs
                new_slide = prs.slides.add_slide(
                    prs.slides[0].slide_layout  # 用第1页同布局(=Master2)，继承MAKE it POSSIBLE装饰
                )
                # 清空布局占位符（不清背景，保留母版装饰）
                for shape in list(new_slide.shapes):
                    shape._element.getparent().remove(shape._element)
                # 复制 tmp_slide 的形状
                for shape in tmp_slide.shapes:
                    self._copy_shape_to_slide(shape, tmp_slide, new_slide, prs)
                self._add_decoration(new_slide)  # 嵌入装饰图到幻灯片本身

            now = datetime.now().strftime('%Y%m%d_%H%M%S')
            out_name = f"批量报告_{now}.pptx"
            out_path = os.path.join(OUTPUT_FOLDER, out_name)
            prs.save(out_path)

            _open_file(out_path)

            self.batch_status_label.config(
                text=f"[生成完成] {out_name}  共{len(records)}页", fg="green")
            self.root.after(4000, lambda: self.batch_status_label.config(text="就绪"))

        except Exception as e:
            messagebox.showerror("生成失败", f"批量生成出错：\n{str(e)}")

    # ====== PPT生成核心逻辑 ======

    def _update_slide_content(self, slide, record):
        photo_path = os.path.join(PHOTO_FOLDER, record['photo'])

        category_clean = record['category'].replace('店铺问题-', '').replace('店铺问题—', '')
        title_text = f"店铺问题-{category_clean}：{record['title']}"

        severity = record.get('severity', '一般问题')
        desc_text = record.get('description', '')
        desc_full = f"[{severity}] {desc_text}"

        remedy_text = record.get('remedy', '')

        ellipse_shapes = []
        photo_left = None
        photo_right = None

        for shape in list(slide.shapes):
            if '椭圆' in shape.name:
                ellipse_shapes.append(shape)
            elif shape.has_text_frame and '[现场照片]' in shape.text_frame.text:
                photo_right = shape
            elif shape.name == 'Rectangle 2097155':
                photo_left = shape

        ellipse_shapes.sort(key=lambda s: s.left)
        photo_rects = [photo_left, photo_right]

        for i, ellipse in enumerate(ellipse_shapes):
            if i < len(photo_rects) and photo_rects[i]:
                photo = photo_rects[i]
                center_x = photo.left + photo.width // 2
                center_y = photo.top + photo.height // 2
                circle_size = int(photo.width * 0.125)
                ellipse.width = circle_size
                ellipse.height = circle_size
                ellipse.left = center_x - circle_size // 2
                ellipse.top = center_y - circle_size // 2

        for shape in list(slide.shapes):
            if shape.has_text_frame and '[现场照片]' in shape.text:
                if os.path.exists(photo_path):
                    try:
                        sp = shape._element
                        sp.getparent().remove(sp)
                        slide.shapes.add_picture(
                            photo_path, shape.left, shape.top,
                            width=shape.width, height=shape.height
                        )
                    except Exception as e:
                        print(f'替换现场照片失败: {e}')
            elif shape.has_text_frame and '[平面索引]' in shape.text:
                plan_img = record.get('plan_img', '')
                if plan_img and plan_img != '无':
                    plan_path = os.path.join(PLAN_IMAGES_FOLDER, plan_img)
                    if os.path.exists(plan_path):
                        try:
                            sp = shape._element
                            sp.getparent().remove(sp)
                            slide.shapes.add_picture(
                                plan_path, shape.left, shape.top,
                                width=shape.width, height=shape.height
                            )
                        except Exception as e:
                            print(f'替换平面索引失败: {e}')

        for shape in ellipse_shapes:
            shape._element.getparent().append(shape._element)

        for shape in list(slide.shapes):
            if shape.has_text_frame:
                text = shape.text_frame.text
                if "店铺问题" in text:
                    self._update_title_text(shape, title_text)
                elif "[问题描述]" in text:
                    self._update_desc_text(shape, desc_full)
                elif "[整改建议]" in text:
                    self._update_remedy_text(shape, remedy_text)

    def _update_title_text(self, shape, new_text):
        tf = shape.text_frame
        parts = new_text.split('：', 1)
        if len(parts) > 1:
            category_part = parts[0].replace('店铺问题-', '')
            title = parts[1]
        else:
            category_part = '分类'
            title = new_text

        for para in tf.paragraphs:
            para.clear()
            run1 = para.add_run()
            run1.text = f"店铺问题-{category_part}："
            run1.font.size = Pt(28)
            run1.font.bold = True
            run1.font.color.rgb = RGBColor(255, 0, 0)

            run2 = para.add_run()
            run2.text = title
            run2.font.size = Pt(28)
            run2.font.bold = True
            run2.font.color.rgb = RGBColor(255, 0, 0)

    def _update_desc_text(self, shape, new_desc):
        tf = shape.text_frame
        if tf.paragraphs:
            para = tf.paragraphs[0]
            para.clear()
            run = para.add_run()
            run.text = new_desc if new_desc else '[问题描述]'
            run.font.size = Pt(12)
            run.font.name = CN_FONT
            run.font.color.rgb = RGBColor(255, 0, 0)
            for i in range(1, len(tf.paragraphs)):
                tf.paragraphs[i].clear()

    def _update_remedy_text(self, shape, new_remedy):
        tf = shape.text_frame
        for para in tf.paragraphs:
            para.clear()
            run = para.add_run()
            run.text = new_remedy if new_remedy else '[整改建议]'
            run.font.size = Pt(12)
            run.font.name = CN_FONT
            run.font.color.rgb = RGBColor(255, 0, 0)

    def _add_decoration(self, slide):
        """在幻灯片最底层添加装饰图（MAKE it POSSIBLE + 红色弧形），保证复制到别的PPT不丢"""
        if not os.path.exists(DECOR_IMG):
            return
        try:
            from pptx.util import Emu
            pic = slide.shapes.add_picture(
                DECOR_IMG,
                Emu(90835), Emu(-5397),
                width=Emu(12232005), height=Emu(6868795)
            )
            # 移到最底层（spTree的第3个位置：nvGrpSpPr, grpSpPr之后）
            spTree = slide.shapes._spTree
            spTree.remove(pic._element)
            spTree.insert(2, pic._element)
        except Exception as e:
            print(f"添加装饰图失败：{e}")

    def _copy_shape_to_slide(self, shape, src_slide, dst_slide, dst_prs):
        try:
            import io
            if shape.shape_type == 13:
                img_bytes = shape.image.blob
                dst_slide.shapes.add_picture(
                    io.BytesIO(img_bytes),
                    shape.left, shape.top,
                    width=shape.width, height=shape.height
                )
            else:
                new_elem = copy.deepcopy(shape._element)
                dst_slide.shapes._spTree.append(new_elem)
        except Exception as e:
            print(f"复制形状失败：{e}")

    # ====== 界面布局 ======

    def create_widgets(self):
        # 顶层 Notebook（选项卡）
        nb = ttk.Notebook(self.root)
        nb.pack(fill=tk.BOTH, expand=True, padx=6, pady=6)

        # ── Tab 1: 单张模式 ──
        tab_single = ttk.Frame(nb)
        nb.add(tab_single, text="  单张录入  ")
        self._build_single_tab(tab_single)

    def _build_single_tab(self, parent):
        main_frame = ttk.Frame(parent, padding="10")
        main_frame.pack(fill=tk.BOTH, expand=True)

        # 顶部标题栏
        top_frame = ttk.Frame(main_frame)
        top_frame.pack(fill=tk.X, pady=(0, 10))

        self.title_label = ttk.Label(top_frame, text="请选择照片", font=(CN_FONT, 14))
        self.title_label.pack(side=tk.LEFT)

        ttk.Label(top_frame, text="华为门店验收工具 v2.8",
                  font=(CN_FONT, 10)).pack(side=tk.RIGHT)

        # 中间：左侧表单 + 右侧照片选择
        middle_frame = ttk.Frame(main_frame)
        middle_frame.pack(fill=tk.BOTH, expand=True)

        # 左侧：表单（缩小宽度）
        form_frame = ttk.LabelFrame(middle_frame, text="问题信息", padding="6")
        form_frame.pack(side=tk.LEFT, fill=tk.BOTH, padx=(0, 10))

        row = 0

        ttk.Label(form_frame, text="类别：", font=(CN_FONT, 9)).grid(row=row, column=0, sticky=tk.W, pady=3)
        default_category = self.settings.get('category', CATEGORIES[0])
        self.category_var = tk.StringVar(value=default_category)
        category_combo = ttk.Combobox(form_frame, textvariable=self.category_var,
                                      values=CATEGORIES, state='readonly', width=14, font=(CN_FONT, 9))
        category_combo.grid(row=row, column=1, sticky=tk.W, pady=3)
        category_combo.bind('<<ComboboxSelected>>', self._on_category_changed)
        row += 1

        ttk.Label(form_frame, text="标题：", font=(CN_FONT, 9)).grid(row=row, column=0, sticky=tk.W, pady=3)
        default_title = self.settings.get('title', '')
        titles_for_category = TITLES_BY_CATEGORY.get(default_category, TITLES_BY_CATEGORY[CATEGORIES[0]])
        self.title_var = tk.StringVar(
            value=default_title if default_title in titles_for_category else titles_for_category[0])
        self.title_combo = ttk.Combobox(form_frame, textvariable=self.title_var,
                                        values=titles_for_category, width=14, font=(CN_FONT, 9))
        self.title_combo.grid(row=row, column=1, sticky=tk.W, pady=3)
        self.title_combo.bind('<<ComboboxSelected>>', self._on_title_changed)
        self.title_combo.bind('<Button-3>', self._title_right_click_batch)
        row += 1

        ttk.Label(form_frame, text="严重程度：", font=(CN_FONT, 9)).grid(row=row, column=0, sticky=tk.W, pady=3)
        default_severity = self.settings.get('severity', SEVERITIES[0])
        self.severity_var = tk.StringVar(value=default_severity)
        severity_combo = ttk.Combobox(form_frame, textvariable=self.severity_var,
                                      values=SEVERITIES, state='readonly', width=14, font=(CN_FONT, 9))
        severity_combo.grid(row=row, column=1, sticky=tk.W, pady=3)
        severity_combo.bind('<<ComboboxSelected>>', self._on_severity_changed)
        row += 1

        ttk.Label(form_frame, text="问题描述：", font=(CN_FONT, 9)).grid(row=row, column=0, sticky=tk.NW, pady=3)
        self.description_text = tk.Text(form_frame, width=22, height=3, font=(CN_FONT, 9))
        self.description_text.grid(row=row, column=1, sticky=tk.W, pady=3)

        db_frame = ttk.Frame(form_frame)
        db_frame.grid(row=row, column=2, sticky=tk.W, padx=3)
        ttk.Button(db_frame, text="从库选", command=self.select_from_issues, width=7).pack(pady=1)
        ttk.Button(db_frame, text="添加", command=self.add_to_issues, width=7).pack(pady=1)
        row += 1

        ttk.Label(form_frame, text="整改意见：", font=(CN_FONT, 9)).grid(row=row, column=0, sticky=tk.NW, pady=3)
        self.remedy_text = tk.Text(form_frame, width=22, height=3, font=(CN_FONT, 9))
        self.remedy_text.grid(row=row, column=1, sticky=tk.W, pady=3)

        remedy_db_frame = ttk.Frame(form_frame)
        remedy_db_frame.grid(row=row, column=2, sticky=tk.W, padx=3)
        ttk.Button(remedy_db_frame, text="从库选", command=self.select_from_remedies, width=7).pack(pady=1)
        ttk.Button(remedy_db_frame, text="添加", command=self.add_to_remedies, width=7).pack(pady=1)
        row += 1

        ttk.Label(form_frame, text="平面图：", font=(CN_FONT, 9)).grid(row=row, column=0, sticky=tk.W, pady=3)
        self.plan_var = tk.StringVar(value="无")
        ttk.Combobox(form_frame, textvariable=self.plan_var,
                     values=self.plan_images, state='readonly', width=14, font=(CN_FONT, 9)).grid(
            row=row, column=1, sticky=tk.W, pady=3)
        row += 1

        nav_frame = ttk.Frame(form_frame)
        nav_frame.grid(row=row, column=0, columnspan=3, pady=10)

        self.record_count_label = tk.Label(nav_frame, text="已记录: 0 张", fg="#FF6600",
                                            font=(CN_FONT, 9, 'bold'))
        self.record_count_label.pack(side=tk.LEFT, padx=(0, 15))

        record_btn = tk.Button(nav_frame, text="记录", command=self.record_current,
                              font=(CN_FONT, 10), bg='#FF9800', fg='white', cursor='hand2')
        record_btn.pack(side=tk.LEFT, padx=(0, 10))

        save_btn = tk.Button(nav_frame, text="生成PPT", command=self.save_current,
                             font=(CN_FONT, 10), bg='#4CAF50', fg='white', cursor='hand2')
        save_btn.pack(side=tk.LEFT)

        self.stats_label = ttk.Label(form_frame, text=f"已处理: {len(self.processed)}/{len(self.photos)} 张", font=(CN_FONT, 8))
        self.stats_label.grid(row=row+1, column=0, columnspan=2, sticky=tk.W, pady=5)

        # ── 已记录列表 ──
        row += 2
        rec_list_lf = ttk.LabelFrame(form_frame, text="已记录列表（单击回看 / Del删除）", padding="4")
        rec_list_lf.grid(row=row, column=0, columnspan=3, sticky='ew', pady=(6, 0))
        form_frame.columnconfigure(0, weight=1)

        rec_list_vsb = ttk.Scrollbar(rec_list_lf, orient=tk.VERTICAL)
        rec_list_vsb.pack(side=tk.RIGHT, fill=tk.Y)
        self._rec_listbox = tk.Listbox(rec_list_lf, height=6, font=(CN_FONT, 9),
                                       yscrollcommand=rec_list_vsb.set,
                                       selectmode=tk.SINGLE, activestyle='dotbox',
                                       exportselection=False)
        self._rec_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        rec_list_vsb.config(command=self._rec_listbox.yview)

        def _on_rec_select(e=None):
            sel = self._rec_listbox.curselection()
            if not sel:
                return
            idx = sel[0]
            if idx < len(self.recorded_data):
                rec = self.recorded_data[idx]
                self._fill_form_with_record(rec)

        def _on_rec_delete(e=None):
            sel = self._rec_listbox.curselection()
            if not sel:
                return
            idx = sel[0]
            if idx >= len(self.recorded_data):
                return
            rec = self.recorded_data[idx]
            fname = rec.get('photo', '')
            if not messagebox.askyesno("确认删除", f"从列表移除第 {idx+1} 条记录？\n{fname}", parent=self.root):
                return
            self.recorded_data.pop(idx)
            # 如果该照片不再有记录，也从 processed 移除
            if fname and not any(r['photo'] == fname for r in self.recorded_data):
                self.processed.discard(fname)
                self._refresh_thumbnail(fname)
            self._refresh_rec_listbox()
            self.record_count_label.config(text=f"已记录: {len(self.recorded_data)} 张")
            self.stats_label.config(text=f"已处理: {len(self.processed)}/{len(self.photos)} 张")

        self._rec_listbox.bind('<<ListboxSelect>>', _on_rec_select)
        self._rec_listbox.bind('<Delete>', _on_rec_delete)

        del_btn_frame = ttk.Frame(rec_list_lf)
        del_btn_frame.pack(side=tk.BOTTOM, fill=tk.X, pady=(3, 0))
        ttk.Button(del_btn_frame, text="✕ 删除选中", command=_on_rec_delete).pack(side=tk.LEFT)

        # 右侧：照片选择区域
        right_frame = ttk.Frame(middle_frame)
        right_frame.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True)

        # 工具栏
        toolbar = ttk.Frame(right_frame)
        toolbar.pack(fill=tk.X, pady=(0, 5))

        ttk.Button(toolbar, text="浏览...",
                    command=self._browse_photo).pack(side=tk.LEFT, padx=(0, 6))

        self._single_status_var = tk.StringVar(value="加载中...")
        ttk.Label(toolbar, textvariable=self._single_status_var,
                   font=(CN_FONT, 8), foreground='#666').pack(side=tk.LEFT, padx=12)

        # 当前选中照片名称
        self.photo_name_label = ttk.Label(right_frame, text="未选择照片",
                                           font=(CN_FONT, 9), foreground='#666')
        self.photo_name_label.pack(pady=(0, 5))

        # 缩略图网格区域（占满剩余空间）
        grid_outer = ttk.Frame(right_frame)
        grid_outer.pack(fill=tk.BOTH, expand=True)

        vsb = ttk.Scrollbar(grid_outer, orient=tk.VERTICAL)
        vsb.pack(side=tk.RIGHT, fill=tk.Y)

        canvas = tk.Canvas(grid_outer, yscrollcommand=vsb.set, highlightthickness=0)
        canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        vsb.config(command=canvas.yview)

        inner = ttk.Frame(canvas)
        canvas_window = canvas.create_window((0, 0), window=inner, anchor='nw')
        self._single_pick_canvas = canvas
        self._single_pick_inner = inner
        self._single_pick_canvas_window = canvas_window

        def _on_canvas_resize(e):
            canvas.itemconfig(canvas_window, width=e.width)
        canvas.bind('<Configure>', _on_canvas_resize)

        def _on_inner_resize(e):
            canvas.configure(scrollregion=canvas.bbox("all"))
        inner.bind('<Configure>', _on_inner_resize)

        # 鼠标滚轮
        def _on_mousewheel(e):
            canvas.yview_scroll(-1 * (e.delta // 120), "units")
        canvas.bind('<Enter>', lambda e: canvas.bind_all('<MouseWheel>', _on_mousewheel))
        canvas.bind('<Leave>', lambda e: canvas.unbind_all('<MouseWheel>'))

        # 底部状态栏
        bottom_frame = ttk.Frame(main_frame)
        bottom_frame.pack(fill=tk.X, pady=(10, 0))

        self.status_label = tk.Label(bottom_frame, text="就绪", fg="blue")
        self.status_label.pack(side=tk.LEFT)

        # 右侧作者信息
        author_label = tk.Label(bottom_frame,
                                text="有改进建议请联系：王柱勇 13945079480",
                                font=(CN_FONT, 8), fg="#888")
        author_label.pack(side=tk.RIGHT)

        # 加载缩略图
        self._load_single_thumb_grid()

# ====== 启动 ======

if __name__ == "__main__":
    os.makedirs(OUTPUT_FOLDER, exist_ok=True)
    os.makedirs(PHOTO_FOLDER, exist_ok=True)

    if not os.path.exists(ISSUES_DB_FILE):
        save_json(ISSUES_DB_FILE, {})

    if not os.path.exists(REMEDIES_DB_FILE):
        save_json(REMEDIES_DB_FILE, {})

    root = tk.Tk()
    app = PhotoReportApp(root)
    root.mainloop()
